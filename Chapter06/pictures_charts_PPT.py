__author__ = 'Chetan'

from pptx import Presentation
from pptx.util import Inches

img_path = 'python.png'
img_path2 = 'learn_python.jpeg'
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(2)
pic = slide.shapes.add_picture(img_path, left, top, height=Inches(2), width=Inches(3))

left = Inches(2)
top = Inches(5)
height = Inches(2)
pic = slide.shapes.add_picture(img_path2, left, top, height=height)

prs.save('picture.pptx')

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Data based on regions'

chart_data = ChartData()
chart_data.categories = ['West', 'East', 'North', 'South']
chart_data.add_series('Series 1', (0.35, 0.25, 0.25, 0.15))

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0%'
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

prs.save('chart.pptx')