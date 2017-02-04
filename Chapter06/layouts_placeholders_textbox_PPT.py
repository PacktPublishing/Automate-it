__author__ = 'Chetan'
# from pptx import Presentation
#
# prs = Presentation()
# two_content_slide_layout = prs.slide_layouts[3]
# slide = prs.slides.add_slide(two_content_slide_layout)
# shapes = slide.shapes
#
# title_shape = shapes.title
# body_shape = shapes.placeholders[1]
# title_shape.text = 'Adding a Two Content Slide'
#
# tf = body_shape.text_frame
# tf.text = 'This is line 1.'
#
# p = tf.add_paragraph()
# p.text = 'Again a Line 2..'
# p.level = 1
#
# p = tf.add_paragraph()
# p.text = 'And this is line 3...'
# p.level = 2
#
# prs.save('two_content.pptx')

from pptx import Presentation
from pptx.util import Inches, Pt
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

txBox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(5), Inches(1))
tf = txBox.text_frame

tf.text = "Wow! I'm inside a textbox"

p = tf.add_paragraph()
p.text = "Adding a new text"
p.font.bold = True
p.font.italic = True
p.font.size = Pt(30)

prs.save('textBox.pptx')

