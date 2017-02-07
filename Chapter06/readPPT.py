#Create presentation object
from pptx import Presentation
path_to_presentation = 'myprofile.pptx'
prs = Presentation(path_to_presentation)
print "Presentation object for myprofile file: \n", prs

#Get presentation objects
print "Presentation Object:", prs
print "Slides are:"
for slide in prs.slides:
    print "Slide object:", slide

#Get slide objects
print "Slide has following objects:"
slide1, slide2 = prs.slides[0], prs.slides[1]
print "Slide Ids: \n", slide1.slide_id, ",", slide2.slide_id
print "Slide Open XML elements: \n", slide1.element, ",", slide2.element
print "Slide layouts: \n", slide1.slide_layout.name, ",", slide2.slide_layout.name

print "Shapes in the slides"
i=1
for slide in prs.slides:
    print 'Slide', i
    for shape in slide.shapes:
        print "Shape: ", shape.shape_type
    i +=1


#Add text to shapes
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)

print "Text is: ", text_runs

