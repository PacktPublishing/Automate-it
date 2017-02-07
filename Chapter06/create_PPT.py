from pptx import Presentation

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Yo, Python!"
slide.placeholders[1].text = "Yes it is really awesome"

prs.save('yoPython.pptx')

from pptx import Presentation
prs = Presentation('sample_ppt.pptx')
first_slide = prs.slides[0]
first_slide.shapes[0].text_frame.paragraphs[0].text = "Hello!"
slide = prs.slides.add_slide(prs.slide_layouts[1])
text_frame = slide.shapes[0].text_frame
p = text_frame.paragraphs[0]
p.text = "This is a paragraph"
prs.save('new_ppt.pptx')




